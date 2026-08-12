from __future__ import annotations

import io
from datetime import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.analysis.models.analysis import (
    AnalysisResult,
    ManagementAction,
)

BANK_NAME = "Meridian Bank"
DECK_TITLE = "Operational Risk Board Review"
CLASSIFICATION = "Internal — Confidential"
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
GREY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

_MONTHS = [
    "January", "February", "March", "April", "May", "June",
    "July", "August", "September", "October", "November", "December",
]
_SLIDE_W = Inches(10)
_SLIDE_H = Inches(7.5)


def _format_month(value: str | None) -> str:
    if not value:
        return ""
    try:
        year, month = value.split("-", 1)
        return f"{_MONTHS[int(month) - 1]} {year}"
    except (ValueError, IndexError):
        return value


def _format_date(value: str | None) -> str:
    if not value:
        return ""
    try:
        dt = datetime.fromisoformat(value.replace("Z", "+00:00"))
    except ValueError:
        return value
    return dt.strftime("%b %d, %Y")


def _period(result: AnalysisResult) -> str:
    if result.trend:
        return _format_month(result.trend[-1].month)
    if result.summary.generatedAt:
        return _format_date(result.summary.generatedAt)
    return _format_date(result.createdAt)


def _metric_by_id(result: AnalysisResult, metric_id: str) -> int | None:
    for metric in result.metrics:
        if metric.id == metric_id:
            return metric.value
    return None


def _clean(value: str | None, limit: int = 220) -> str:
    text = (value or "").replace("\n", " ").replace("\r", " ").strip()
    if len(text) > limit:
        return text[: limit - 1].rstrip() + "…"
    return text


def _add_slide(prs: Presentation, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    title_shape = slide.shapes.title
    title_shape.text = title
    for run in title_shape.text_frame.paragraphs[0].runs:
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = NAVY
    # Accent bar under the title.
    bar = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(1.4), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    return slide


def _add_textbox(slide, left: float, top: float, width: float, height: float):
    return slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )


def _add_bullets(slide, items: list[str], top: float = 1.5, size: int = 14) -> None:
    box = _add_textbox(slide, 0.6, top, 8.8, 5.2)
    frame = box.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = NAVY if index == 0 else RGBColor(0x33, 0x37, 0x41)
    box.text_frame.paragraphs[0].font.size = Pt(size)


def _add_table(slide, headers: list[str], rows: list[list[str]], top: float, height: float = 0.5) -> None:
    shape = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(0.6),
        Inches(top),
        Inches(8.8),
        Inches(height),
    )
    table = shape.table
    table.columns[0].width = Inches(0.9)
    for index, header in enumerate(headers):
        cell = table.cell(0, index)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(11)
                run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = _clean(value)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
            if row_index % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xEF, 0xF3, 0xFA)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE


def _footer(slide, left: str, right: str) -> None:
    box = _add_textbox(slide, 0.6, 7.05, 8.8, 0.35)
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = f"{left}   ·   {right}"
    paragraph.alignment = PP_ALIGN.CENTER
    for run in paragraph.runs:
        run.font.size = Pt(9)
        run.font.color.rgb = GREY


def build_powerpoint_deck(result: AnalysisResult) -> bytes:
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H

    period = _period(result)
    generated = _format_date(result.summary.generatedAt or result.createdAt)

    # Slide 1 — Title --------------------------------------------------------------
    slide = _add_slide(prs, BANK_NAME)
    box = _add_textbox(slide, 0.6, 1.9, 8.8, 3.6)
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = DECK_TITLE
    paragraph.alignment = PP_ALIGN.CENTER
    for run in paragraph.runs:
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = NAVY
    paragraph = frame.add_paragraph()
    paragraph.text = f"Monthly Risk Position — {period}"
    paragraph.alignment = PP_ALIGN.CENTER
    for run in paragraph.runs:
        run.font.size = Pt(20)
        run.font.color.rgb = GREY
    paragraph = frame.add_paragraph()
    paragraph.text = CLASSIFICATION
    paragraph.alignment = PP_ALIGN.CENTER
    for run in paragraph.runs:
        run.font.size = Pt(12)
        run.font.color.rgb = ACCENT
    _footer(slide, f"Prepared by Risk Copilot AI · {generated}", CLASSIFICATION)

    # Slide 2 — Executive Summary ---------------------------------------------------
    slide = _add_slide(prs, "Executive Summary")
    score = result.overallScore
    headline = f"Overall risk score {score.score}/100 — {score.level}"
    paragraphs = result.summary.paragraphs or [
        "No executive summary was generated for this analysis run."
    ]
    body = [_clean(paragraphs[0], 300)]
    if result.metrics:
        body.append("Headline metrics: " + ", ".join(
            f"{metric.label} {metric.value}" for metric in result.metrics
        ))
    _add_bullets(slide, [headline, *body])
    _footer(slide, f"Generated by Risk Copilot AI · {generated}", CLASSIFICATION)

    # Slide 3 — Risk Landscape ------------------------------------------------------
    slide = _add_slide(prs, "Risk Landscape")
    description = score.description or "No description available."
    _add_bullets(slide, [
        f"Overall score {score.score}/100 ({score.level})"
        + (f" — {score.change}" if score.change else ""),
        description,
    ], top=1.5, size=13)
    distribution = [
        ("Critical", _metric_by_id(result, "critical-risks")),
        ("High", _metric_by_id(result, "high-risks")),
        ("Medium", _metric_by_id(result, "medium-risks")),
        ("Low", _metric_by_id(result, "low-risks")),
    ]
    _add_table(
        slide,
        ["Severity", "Risk Count"],
        [[severity, str(value) if value is not None else "n/a"] for severity, value in distribution],
        top=3.0,
    )
    if result.heatmap:
        exposure = [
            f"{row.division}: " + ", ".join(f"{cell.category} ({cell.level})" for cell in row.cells[:3])
            for row in result.heatmap[:4]
        ]
        if exposure:
            box = _add_textbox(slide, 0.6, 4.7, 8.8, 2.2)
            frame = box.text_frame
            frame.word_wrap = True
            frame.paragraphs[0].text = "Concentration by division:"
            frame.paragraphs[0].runs[0].font.bold = True
            frame.paragraphs[0].runs[0].font.size = Pt(12)
            for item in exposure:
                paragraph = frame.add_paragraph()
                paragraph.text = item
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.color.rgb = GREY
    _footer(slide, f"Risk Copilot AI · {period}", CLASSIFICATION)

    # Slide 4 — Top Risks -----------------------------------------------------------
    slide = _add_slide(prs, "Top Risks")
    findings = result.keyFindings or []
    if not findings:
        _add_bullets(slide, ["No key findings were generated for this analysis run."])
    else:
        rows = [
            [
                str(index + 1),
                finding.id,
                finding.severity,
                _clean(finding.title, 100),
                finding.category,
                finding.likelihood or "",
                finding.exposure or "",
            ]
            for index, finding in enumerate(findings[:5])
        ]
        _add_table(
            slide,
            ["Rank", "ID", "Severity", "Risk", "Category", "Likelihood", "Exposure"],
            rows,
            top=1.6,
            height=0.9,
        )
    _footer(slide, f"Source evidence cited in the full report · {generated}", CLASSIFICATION)

    # Slide 5 — Audit Findings ------------------------------------------------------
    slide = _add_slide(prs, "Audit Findings")
    audit_rows = [
        ("Critical findings", _metric_by_id(result, "critical-findings")),
        ("Open audit findings", _metric_by_id(result, "open-findings")),
        ("Overdue audit findings", _metric_by_id(result, "overdue-findings")),
    ]
    _add_table(
        slide,
        ["Audit metric", "Value"],
        [[label, str(value) if value is not None else "n/a"] for label, value in audit_rows],
        top=1.8,
    )
    audit_docs = [doc for doc in result.documents if doc.category in ("audit-findings", "gia-findings")]
    note = (
        "Aggregate audit figures are derived from "
        + ", ".join(doc.filename for doc in audit_docs)
        + "."
        if audit_docs
        else "No audit finding documents were included in this analysis run."
    )
    _add_bullets(slide, [note], top=3.4, size=12)
    _footer(slide, f"Risk Copilot AI · {generated}", CLASSIFICATION)

    # Slide 6 — Management Actions --------------------------------------------------
    slide = _add_slide(prs, "Management Actions")
    actions = result.managementActions or []
    if not actions:
        _add_bullets(slide, ["No management actions were recorded for this analysis run."])
    else:
        rows = [
            [
                action.id,
                action.priority,
                _clean(action.action, 90),
                action.owner,
                action.department,
                action.dueDate,
                action.status,
            ]
            for action in actions[:6]
        ]
        _add_table(
            slide,
            ["ID", "Priority", "Action", "Owner", "Department", "Due Date", "Status"],
            rows,
            top=1.6,
            height=1.0,
        )
    _footer(slide, f"Full detail in the Monthly Operational Risk Report · {period}", CLASSIFICATION)

    # Slide 7 — Closing / Next Steps -------------------------------------------------
    slide = _add_slide(prs, "Next Steps")
    box = _add_textbox(slide, 0.6, 1.7, 8.8, 1.1)
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = "Questions & Discussion"
    paragraph.alignment = PP_ALIGN.CENTER
    for run in paragraph.runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = NAVY
    recommendations = result.recommendations or []
    next_steps: list[str] = []
    if recommendations:
        for recommendation in recommendations[:4]:
            next_steps.append(f"[{recommendation.priority}] {_clean(recommendation.action, 140)}")
    else:
        next_steps.append("No AI recommendations were generated for this analysis run.")
    next_steps.append("Full audit-level detail in the Monthly Operational Risk Report.")
    _add_bullets(slide, next_steps, top=3.2, size=14)
    box = _add_textbox(slide, 0.6, 6.3, 8.8, 0.6)
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = (
        "This deck was generated by Risk Copilot AI from uploaded source documents "
        "and is for internal management use only."
    )
    paragraph.alignment = PP_ALIGN.CENTER
    for run in paragraph.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = GREY
    _footer(slide, f"{BANK_NAME} · Risk Copilot AI", CLASSIFICATION)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
