import io
import json

import pytest
from docx import Document
from fastapi.testclient import TestClient
from pptx import Presentation

from app.analysis.models import AnalysisResult
from app.analysis.models.analysis import (
    AnalysisMetric,
    AnalysisSummary,
    DocumentCoverage,
    Evidence,
    HeatmapCell,
    HeatmapRow,
    KeyFinding,
    ManagementAction,
    OverallScore,
    Recommendation,
    RiskTrendPoint,
    SourceCount,
)
from app.main import app
from app.services import analysis_store

client = TestClient(app)

DOCX_MEDIA = (
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
)
PPTX_MEDIA = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)


@pytest.fixture(autouse=True)
def _isolate_store(tmp_path, monkeypatch):
    monkeypatch.setattr(analysis_store, "RUNS_DIR", tmp_path / "runs")
    monkeypatch.setattr(analysis_store, "RESULTS_DIR", tmp_path / "results")


def _sample_result() -> AnalysisResult:
    return AnalysisResult(
        id="res-export-1",
        status="ready",
        createdAt="2026-08-06T10:00:00Z",
        confidence=100,
        warnings=[
            "AI synthesis unavailable; executive narrative generated deterministically."
        ],
        documents=[
            DocumentCoverage(
                id="d1",
                filename="Risk Register.xlsx",
                category="risk-register",
                status="interpreted",
                evidenceCount=38,
            )
        ],
        overallScore=OverallScore(
            score=64,
            level="High",
            description="Elevated operational risk exposure across the bank.",
            change="-6 points vs last month",
        ),
        metrics=[
            AnalysisMetric(id="critical-risks", label="Critical Risks", value=4),
            AnalysisMetric(id="high-risks", label="High Risks", value=15),
            AnalysisMetric(id="medium-risks", label="Medium Risks", value=14),
            AnalysisMetric(id="low-risks", label="Low Risks", value=5),
            AnalysisMetric(id="total-risks", label="Total Risks", value=38),
            AnalysisMetric(id="critical-findings", label="Critical Findings", value=3),
            AnalysisMetric(id="open-findings", label="Open Audit Findings", value=18),
            AnalysisMetric(id="overdue-findings", label="Overdue Audit Findings", value=3),
            AnalysisMetric(id="open-exceptions", label="Open Exceptions", value=12),
            AnalysisMetric(id="overdue-exceptions", label="Overdue Exceptions", value=8),
        ],
        heatmap=[
            HeatmapRow(
                division="Technology",
                cells=[
                    HeatmapCell(category="Cyber", level=3),
                    HeatmapCell(category="People", level=2),
                ],
            )
        ],
        trend=[
            RiskTrendPoint(month="2026-06", high=21, medium=13, low=4),
            RiskTrendPoint(month="2026-07", high=19, medium=14, low=5),
        ],
        keyFindings=[
            KeyFinding(
                id="RISK-2026-001",
                title="Privileged access not recertified",
                category="Cyber",
                severity="Critical",
                likelihood="Very likely",
                exposure="High",
                evidence=[
                    Evidence(
                        documentId="d1",
                        documentType="risk-register",
                        sourceRef="Risk Register · row 2",
                        snippet="privileged access not recertified for 8+ months",
                    )
                ],
                confidence=92,
            )
        ],
        recommendations=[
            Recommendation(
                id="rec-export-1",
                priority="Critical",
                category="Identity & Access",
                action="Recertify privileged access immediately.",
                impact="Reduces critical exposure by ~30%",
                evidence=[
                    Evidence(
                        documentId="d1",
                        documentType="risk-register",
                        sourceRef="Risk Register · row 2",
                    )
                ],
                confidence=88,
            )
        ],
        managementActions=[
            ManagementAction(
                id="ACT-EXPORT-01",
                action="Recertify legacy system privileged access",
                owner="D. Patel",
                department="Technology",
                dueDate="Aug 15, 2026",
                priority="Critical",
                status="In Progress",
            )
        ],
        summary=AnalysisSummary(
            generatedAt="2026-08-06T10:05:00Z",
            paragraphs=[
                "Meridian Bank is in a high operational risk position.",
                "Critical findings remain concentrated in identity and access management.",
            ],
            sources=[SourceCount(label="Risk Register.xlsx", count="38 rows")],
        ),
    )


def _save(result: AnalysisResult) -> None:
    analysis_store.save_result(result)


def _save_run(result_id: str, run_id: str = "run-export-1") -> None:
    from app.analysis.models import AnalysisRun

    analysis_store.save_run(
        AnalysisRun(id=run_id, documentIds=["d1"], status="ready", resultId=result_id)
    )


def _docx_text(content: bytes) -> str:
    doc = Document(io.BytesIO(content))
    parts = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return "\n".join(parts)


def _pptx_text(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
    return "\n".join(parts)


# --- Word ------------------------------------------------------------------


def test_export_report_returns_valid_docx_with_real_values():
    _save(_sample_result())

    response = client.get("/api/exports/report", params={"result_id": "res-export-1"})
    assert response.status_code == 200
    assert response.headers["content-type"].startswith(DOCX_MEDIA)
    assert "filename=" in response.headers["content-disposition"]
    assert response.headers["content-disposition"].endswith(".docx\"")
    assert response.content[:2] == b"PK"

    text = _docx_text(response.content)
    assert "Meridian Bank" in text
    assert "Internal — Confidential" in text
    assert "64/100" in text
    assert "High" in text
    assert "RISK-2026-001" in text
    assert "Privileged access not recertified" in text
    assert "Recertify privileged access immediately." in text
    assert "ACT-EXPORT-01" in text
    assert "Risk Register · row 2" in text
    assert "Confidence and Disclaimer" in text


def test_export_report_missing_result_returns_404():
    response = client.get("/api/exports/report", params={"result_id": "nope"})
    assert response.status_code == 404


def test_export_report_missing_run_returns_404():
    response = client.get("/api/exports/report", params={"run_id": "nope"})
    assert response.status_code == 404


def test_export_report_without_identifier_returns_422():
    response = client.get("/api/exports/report")
    assert response.status_code == 422


def test_export_report_by_run_id():
    _save(_sample_result())
    _save_run("res-export-1")
    response = client.get("/api/exports/report", params={"run_id": "run-export-1"})
    assert response.status_code == 200
    assert response.content[:2] == b"PK"


def test_export_report_malformed_json_returns_404(tmp_path):
    analysis_store.RESULTS_DIR.mkdir(parents=True, exist_ok=True)
    (analysis_store.RESULTS_DIR / "res-broken-1.json").write_text(
        "{not json", encoding="utf-8"
    )
    response = client.get("/api/exports/report", params={"result_id": "res-broken-1"})
    assert response.status_code == 404


def test_export_report_invalid_model_returns_404(tmp_path):
    analysis_store.RESULTS_DIR.mkdir(parents=True, exist_ok=True)
    (analysis_store.RESULTS_DIR / "res-broken-2.json").write_text(
        json.dumps({"id": 123, "createdAt": []}), encoding="utf-8"
    )
    response = client.get("/api/exports/report", params={"result_id": "res-broken-2"})
    assert response.status_code == 404


def test_export_report_empty_result_still_generates():
    from app.analysis.models.analysis import AnalysisResult

    _save(AnalysisResult(id="res-empty-1", createdAt="2026-08-06T10:00:00Z"))
    response = client.get("/api/exports/report", params={"result_id": "res-empty-1"})
    assert response.status_code == 200
    text = _docx_text(response.content)
    assert "Meridian Bank" in text
    assert "No executive summary was generated" in text


# --- PowerPoint -------------------------------------------------------------


def test_export_presentation_returns_valid_pptx_with_real_values():
    _save(_sample_result())

    response = client.get(
        "/api/exports/presentation", params={"result_id": "res-export-1"}
    )
    assert response.status_code == 200
    assert response.headers["content-type"].startswith(PPTX_MEDIA)
    assert response.headers["content-disposition"].endswith(".pptx\"")
    assert response.content[:2] == b"PK"

    prs = Presentation(io.BytesIO(response.content))
    assert len(list(prs.slides)) == 7
    text = _pptx_text(response.content)
    assert "Meridian Bank" in text
    assert "Executive Summary" in text
    assert "Risk Landscape" in text
    assert "Top Risks" in text
    assert "Audit Findings" in text
    assert "Management Actions" in text
    assert "Next Steps" in text
    assert "64/100" in text
    assert "RISK-2026-001" in text
    assert "ACT-EXPORT-01" in text


def test_export_presentation_missing_result_returns_404():
    response = client.get(
        "/api/exports/presentation", params={"result_id": "nope"}
    )
    assert response.status_code == 404


def test_export_presentation_without_identifier_returns_422():
    response = client.get("/api/exports/presentation")
    assert response.status_code == 422


def test_export_presentation_by_run_id():
    _save(_sample_result())
    _save_run("res-export-1")
    response = client.get(
        "/api/exports/presentation", params={"run_id": "run-export-1"}
    )
    assert response.status_code == 200
    assert response.content[:2] == b"PK"
